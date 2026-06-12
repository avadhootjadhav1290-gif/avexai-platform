# modules/rag_manager.py

from sentence_transformers import SentenceTransformer
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
from io import BytesIO
import re

from db import supabase

# Load embedding model once
import streamlit as st

@st.cache_resource
def get_embedding_model():
    return SentenceTransformer(
        "all-MiniLM-L6-v2"
    )


def extract_pdf_text(file_bytes):

    pdf = PdfReader(
        BytesIO(file_bytes)
    )

    text = ""

    for page in pdf.pages:

        page_text = page.extract_text()

        if page_text:

            page_text = page_text.replace(
                "\x00",
                ""
            )

            text += page_text + "\n"

    return text


def extract_docx_text(file_bytes):

    doc = Document(
        BytesIO(file_bytes)
    )

    text = ""

    for paragraph in doc.paragraphs:

        if paragraph.text:

            text += paragraph.text + "\n"

    return text

def extract_txt_text(file_bytes):

    return file_bytes.decode(
        "utf-8",
        errors="ignore"
    )
    
def extract_pptx_text(file_bytes):

    prs = Presentation(
        BytesIO(file_bytes)
    )

    text = ""

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text += shape.text + "\n"

    return text

def extract_image_text(file_bytes):

    image = Image.open(
        BytesIO(file_bytes)
    )

    return pytesseract.image_to_string(image)

def chunk_text(
    text,
    chunk_size=1000
):

    text = text.replace(
        "\x00",
        ""
    )

    text = re.sub(
        r"[\x00-\x08\x0B\x0C\x0E-\x1F]",
        "",
        text
    )

    chunks = []

    for i in range(
        0,
        len(text),
        chunk_size
    ):

        chunk = text[
            i:i + chunk_size
        ]

        chunk = chunk.replace(
            "\x00",
            ""
        )

        chunks.append(
            chunk
        )

    return chunks


def generate_embedding(text):

    model = get_embedding_model()

    embedding = model.encode(text)

    return embedding.tolist()


def search_similar_chunks(
    question,
    project_id,
    top_k=5
):

    query_embedding = generate_embedding(
        question
    )

    result = supabase.rpc(
        "match_document_chunks",
        {
            "query_embedding": query_embedding,
            "match_count": top_k,
            "project_filter": project_id
        }
    ).execute()

    return result.data


def build_context(
    chunks
):

    context = ""

    for chunk in chunks:

        context += chunk["chunk_text"]
        context += "\n\n"

    return context


def get_document_context(
    question,
    project_id
):

    chunks = search_similar_chunks(
        question,
        project_id,
        top_k=5
    )

    return build_context(
        chunks
    )
    
def search_chunks_by_files(
    question,
    file_ids,
    top_k=5
):

    query_embedding = generate_embedding(
        question
    )

    result = supabase.rpc(
        "match_chat_file_chunks",
        {
            "query_embedding": query_embedding,
            "file_ids": file_ids,
            "match_count": top_k
        }
    ).execute()

    return result.data

def get_chat_file_context(
    question,
    file_ids
):

    if not file_ids:
        return ""

    chunks = search_chunks_by_files(
        question,
        file_ids,
        top_k=5
    )

    return build_context(
        chunks
    )